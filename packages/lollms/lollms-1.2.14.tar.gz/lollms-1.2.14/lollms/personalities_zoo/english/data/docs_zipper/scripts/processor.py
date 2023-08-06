from lollms.personality import APScript, AIPersonality
from lollms.helpers import ASCIIColors

import torch
import numpy as np
import torch
from sklearn.metrics.pairwise import cosine_similarity
from transformers import AutoTokenizer, AutoModel
import json
from pathlib import Path
import PyPDF2
import matplotlib.pyplot as plt
import numpy as np
from sklearn.manifold import TSNE
from sklearn.decomposition import PCA

from sklearn.metrics.pairwise import cosine_similarity
import plotly.graph_objects as go
import seaborn as sns
import matplotlib.pyplot as plt
import mplcursors
import textwrap
from tkinter import Tk, Text, Scrollbar, Frame, Label, TOP, BOTH, RIGHT, LEFT, Y, N, END
from docx import Document
from bs4 import BeautifulSoup
import json
import csv
from pptx import Presentation
from tqdm import tqdm
from pptx.util import Pt
class TextVectorizer:
    def __init__(self, model_name, database_file:Path|str, visualize_data_at_startup=False, visualize_data_at_add_file=False, visualize_data_at_generate=False):
        self.tokenizer = AutoTokenizer.from_pretrained(model_name)
        self.model = AutoModel.from_pretrained(model_name)
        self.embeddings = {}
        self.texts = {}
        self.ready = False
        self.database_file = Path(database_file)
        self.visualize_data_at_startup  = visualize_data_at_startup
        self.visualize_data_at_add_file = visualize_data_at_add_file
        self.visualize_data_at_generate = visualize_data_at_generate

        # Load previous state from the JSON file
        if Path(self.database_file).exists():
            self.load_from_json()
            if visualize_data_at_startup:
                self.show_document()
                
    def show_document(self, query_text="What is the main idea of this text?", use_pca=True):
        if use_pca:
            print("Showing pca representation :")
        else:
            print("Showing t-sne representation :")
        texts = list(self.texts.values())
        embeddings = torch.stack(list(self.embeddings.values())).detach().squeeze(1).numpy()
        # Normalize embeddings
        norms = np.linalg.norm(embeddings, axis=1)
        normalized_embeddings = embeddings / norms[:, np.newaxis]

        # Embed the query text
        query_embedding = self.embed_query(query_text)
        query_embedding = query_embedding.detach().squeeze().numpy()
        query_normalized_embedding = query_embedding / np.linalg.norm(query_embedding)

        # Combine the query embedding with the document embeddings
        combined_embeddings = np.vstack((normalized_embeddings, query_normalized_embedding))

        if use_pca:
            # Use PCA for dimensionality reduction
            pca = PCA(n_components=2)
            embeddings_2d = pca.fit_transform(combined_embeddings)
        else:
            # Use t-SNE for dimensionality reduction
            # Adjust the perplexity value
            perplexity = min(30, combined_embeddings.shape[0] - 1)
            tsne = TSNE(n_components=2, perplexity=perplexity)
            embeddings_2d = tsne.fit_transform(combined_embeddings)


        # Create a scatter plot using Seaborn
        sns.scatterplot(x=embeddings_2d[:-1, 0], y=embeddings_2d[:-1, 1])  # Plot document embeddings
        plt.scatter(embeddings_2d[-1, 0], embeddings_2d[-1, 1], color='red')  # Plot query embedding

        # Add labels to the scatter plot
        for i, (x, y) in enumerate(embeddings_2d[:-1]):
            plt.text(x, y, str(i), fontsize=8)

        plt.xlabel('Dimension 1')
        plt.ylabel('Dimension 2')
        plt.title('Embeddings Scatter Plot based on t-SNE')

        # Enable mplcursors to show tooltips on hover
        cursor = mplcursors.cursor(hover=True)

        # Define the hover event handler
        @cursor.connect("add")
        def on_hover(sel):
            index = sel.target.index
            if index > 0:
                text = texts[index]
                wrapped_text = textwrap.fill(text, width=50)  # Wrap the text into multiple lines
                sel.annotation.set_text(f"Index: {index}\nText:\n{wrapped_text}")
            else:
                sel.annotation.set_text("Query")

        # Define the click event handler using matplotlib event handling mechanism
        def on_click(event):
            if event.xdata is not None and event.ydata is not None:
                x, y = event.xdata, event.ydata
                distances = ((embeddings_2d[:, 0] - x) ** 2 + (embeddings_2d[:, 1] - y) ** 2)
                index = distances.argmin()
                text = texts[index] if index < len(texts) else query_text

                # Open a new Tkinter window with the content of the text
                root = Tk()
                root.title(f"Text for Index {index}")
                frame = Frame(root)
                frame.pack(fill=BOTH, expand=True)

                label = Label(frame, text="Text:")
                label.pack(side=TOP, padx=5, pady=5)

                text_box = Text(frame)
                text_box.pack(side=TOP, padx=5, pady=5, fill=BOTH, expand=True)
                text_box.insert(END, text)

                scrollbar = Scrollbar(frame)
                scrollbar.pack(side=RIGHT, fill=Y)
                scrollbar.config(command=text_box.yview)
                text_box.config(yscrollcommand=scrollbar.set)

                text_box.config(state="disabled")

                root.mainloop()

        # Connect the click event handler to the figure
        plt.gcf().canvas.mpl_connect("button_press_event", on_click)
        plt.show()
        
    def index_document(self, document_id, text, chunk_size, overlap_size, force_vectorize=False):
        if document_id in self.embeddings and not force_vectorize:
            print(f"Document {document_id} already exists. Skipping vectorization.")
            return

        # Tokenize text
        tokens = self.tokenizer.encode_plus(text, add_special_tokens=False, return_attention_mask=False)['input_ids']

        # Split tokens into sentences
        sentences = self.tokenizer.decode(tokens).split('. ')

        # Generate chunks with overlap and sentence boundaries
        chunks = []
        current_chunk = []
        for sentence in sentences:
            sentence_tokens = self.tokenizer.encode_plus(sentence, add_special_tokens=False, return_attention_mask=False)['input_ids']
            if len(current_chunk) + len(sentence_tokens) <= chunk_size:
                current_chunk.extend(sentence_tokens)
            else:
                if current_chunk:
                    chunks.append(current_chunk)
                current_chunk = sentence_tokens

        if current_chunk:
            chunks.append(current_chunk)

        # Generate overlapping chunks
        overlapping_chunks = []
        for i in range(len(chunks)):
            chunk_start = i * (chunk_size - overlap_size)
            chunk_end = min(chunk_start + chunk_size, len(tokens))
            chunk = tokens[chunk_start:chunk_end]
            overlapping_chunks.append(chunk)

        # Generate embeddings for each chunk
        for i, chunk in enumerate(overlapping_chunks):
            # Pad the chunk if it is smaller than chunk_size
            if len(chunk) < chunk_size:
                padding = [self.tokenizer.pad_token_id] * (chunk_size - len(chunk))
                chunk.extend(padding)

            # Convert tokens to IDs
            input_ids = chunk[:chunk_size]

            # Convert input to PyTorch tensor
            input_tensor = torch.tensor([input_ids])

            # Generate chunk embedding
            with torch.no_grad():
                self.model.eval()
                outputs = self.model(input_tensor)
                embeddings = outputs.last_hidden_state.mean(dim=1)

            # Store chunk ID, embedding, and original text
            chunk_id = f"{document_id}_chunk_{i + 1}"
            self.embeddings[chunk_id] = embeddings
            self.texts[chunk_id] = self.tokenizer.decode(chunk[:chunk_size], skip_special_tokens=True)

        self.save_to_json()
        self.ready = True
        if self.visualize_data_at_add_file:
            self.show_document()


    def embed_query(self, query_text):
        # Tokenize query text
        query_tokens = self.tokenizer.encode(query_text)

        # Convert input to PyTorch tensor
        query_input_tensor = torch.tensor([query_tokens])

        # Generate query embedding
        with torch.no_grad():
            self.model.eval()
            query_outputs = self.model(query_input_tensor)
            query_embedding = query_outputs.last_hidden_state.mean(dim=1)

        return query_embedding

    def recover_text(self, query_embedding, top_k=1):
        similarities = {}
        for chunk_id, chunk_embedding in self.embeddings.items():
            similarity = cosine_similarity(query_embedding.numpy(), chunk_embedding.numpy())[0][0]
            similarities[chunk_id] = similarity

        # Sort the similarities and retrieve the top-k most similar embeddings
        sorted_similarities = sorted(similarities.items(), key=lambda x: x[1], reverse=True)[:top_k]

        # Retrieve the original text associated with the most similar embeddings
        texts = [self.texts[chunk_id] for chunk_id, _ in sorted_similarities]

        if self.visualize_data_at_generate:
            self.show_document()

        return texts

    def save_to_json(self):
        state = {
            "embeddings": {str(k): v.tolist() for k, v in self.embeddings.items()},
            "texts": self.texts,
        }
        with open(self.database_file, "w") as f:
            json.dump(state, f)

    def load_from_json(self):
        ASCIIColors.info("Loading vectorized documents")
        with open(self.database_file, "r") as f:
            state = json.load(f)
            self.embeddings = {k: torch.tensor(v) for k, v in state["embeddings"].items()}
            self.texts = state["texts"]
            self.ready = True


class Processor(APScript):
    """
    A class that processes model inputs and outputs.

    Inherits from APScript.
    """

    def __init__(self, personality: AIPersonality, model = None) -> None:
        super().__init__()
        self.personality = personality
        self.model = model
        default_config = {
            "database_path":    f"{self.personality.name}_db.json",
            "chunk_size":       512, # Max chunk size
            "overlap_size":     50, # Overlap between text chunks
            "max_answer_size":  128,  # Maximum answer size
            "visualize_data_at_startup":            False,
            "visualize_data_at_add_file":           False,            
            "visualize_data_at_generate":           False,            
        }
        try:
            self.config = self.load_config_file(self.personality.lollms_paths.personal_configuration_path/f"personality_{self.personality.name}.yaml", default_config)
        except Exception as ex:
            print("Error loading configuration file.\nTrying to reinstall.")
            self.install_personality(Path(__file__).parent.parent, force_reinstall=True)
            self.config = self.load_config_file(self.personality.lollms_paths.personal_configuration_path/f"personality_{self.personality.name}.yaml", default_config)


        self.vector_store = TextVectorizer(
                                    "bert-base-uncased", 
                                    self.personality.lollms_paths.personal_data_path/self.config["database_path"],
                                    visualize_data_at_startup=self.config["visualize_data_at_startup"],
                                    visualize_data_at_add_file=self.config["visualize_data_at_add_file"],
                                    visualize_data_at_generate=self.config["visualize_data_at_generate"]
                                    )
        self.personality.lollms_paths.personal_configuration_path/f"personality_{self.personality.name}.yaml"

    @staticmethod        
    def read_pdf_file(file_path):
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
        return text

    @staticmethod
    def read_docx_file(file_path):
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text

    @staticmethod
    def read_json_file(file_path):
        with open(file_path, 'r') as file:
            data = json.load(file)
        return data
    
    @staticmethod
    def read_csv_file(file_path):
        with open(file_path, 'r') as file:
            csv_reader = csv.reader(file)
            lines = [row for row in csv_reader]
        return lines    

    @staticmethod
    def read_html_file(file_path):
        with open(file_path, 'r') as file:
            soup = BeautifulSoup(file, 'html.parser')
            text = soup.get_text()
        return text
    @staticmethod
    def read_pptx_file(file_path):
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
        return text
    @staticmethod
    def read_text_file(file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        return content
    
    def build_db(self):
        ASCIIColors.info("-> Vectorizing the database"+ASCIIColors.color_orange)
        for file in self.files:
            try:
                if Path(file).suffix==".pdf":
                    text =  Processor.read_pdf_file(file)
                elif Path(file).suffix==".docx":
                    text =  Processor.read_docx_file(file)
                elif Path(file).suffix==".docx":
                    text =  Processor.read_pptx_file(file)
                elif Path(file).suffix==".json":
                    text =  Processor.read_json_file(file)
                elif Path(file).suffix==".csv":
                    text =  Processor.read_csv_file(file)
                elif Path(file).suffix==".html":
                    text =  Processor.read_html_file(file)
                else:
                    text =  Processor.read_text_file(file)
                try:
                    chunk_size=int(self.config["chunk_size"])
                except:
                    ASCIIColors.warning(f"Couldn't read chunk size. Verify your configuration file")
                    chunk_size=512
                try:
                    overlap_size=int(self.config["overlap_size"])
                except:
                    ASCIIColors.warning(f"Couldn't read chunk size. Verify your configuration file")
                    overlap_size=50

                self.vector_store.index_document(file, text, chunk_size=chunk_size, overlap_size=overlap_size)
                
                print(ASCIIColors.color_reset)
                ASCIIColors.success(f"File {file} vectorized successfully")
            except Exception as ex:
                ASCIIColors.error(f"Couldn't vectorize {file}: The vectorizer threw this exception:{ex}")

    def add_file(self, path):
        super().add_file(path)
        try:
            self.build_db()
            return True
        except Exception as ex:
            ASCIIColors.error(f"Couldn't vectorize the database: The vectgorizer threw this exception: {ex}")
            return False        

    

    def run_workflow(self, prompt, previous_discussion_text="", callback=None):
        """
        Runs the workflow for processing the model input and output.

        This method should be called to execute the processing workflow.

        Args:
            generate_fn (function): A function that generates model output based on the input prompt.
                The function should take a single argument (prompt) and return the generated text.
            prompt (str): The input prompt for the model.
            previous_discussion_text (str, optional): The text of the previous discussion. Default is an empty string.

        Returns:
            None
        """
        output =""

        if prompt.strip().lower()=="zip":
            compressed_text = ""
            ASCIIColors.info("Compressing database")
            for chunk in tqdm(self.vector_store.texts):
                prmpt = "Doc: "+chunk+"\nSummary:"
                compressed_text += self.generate(prmpt, self.config["max_answer_size"])
            self.vector_store  = TextVectorizer(
                                    "bert-base-uncased", 
                                    self.personality.lollms_paths.personal_data_path/self.config["database_path"][:-5]+"_compressed.json",
                                    visualize_data_at_startup=self.config["visualize_data_at_startup"],
                                    visualize_data_at_add_file=self.config["visualize_data_at_add_file"],
                                    visualize_data_at_generate=self.config["visualize_data_at_generate"]
                                    )
            self.vector_store.index_document("back", compressed_text, chunk_size=self.config["chunk_size"], overlap_size=self.config["overlap_size"])
            
        elif prompt.strip().lower()=="save":

        docs = self.vector_store.recover_text(self.vector_store.embed_query(prompt), top_k=3)
        docs = '\nDoc:\n'.join([v for v in docs])
        full_text = self.personality.personality_conditioning+"\n### Documentation:\nDoc:\n"+docs+"\n### Question: "+prompt+"\n### Answer:"
        ASCIIColors.error("-------------- Documentation -----------------------")
        ASCIIColors.error(full_text)
        ASCIIColors.error("----------------------------------------------------")
        print("Thinking")
        output = self.generate(full_text, self.config["max_answer_size"])
        return output



